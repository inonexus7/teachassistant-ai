import io
import json
import os

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from urllib.parse import quote_plus
from dotenv import load_dotenv
import datetime
import random
from pyunsplash import PyUnsplash
import re

dir_path = 'static/presentations'

load_dotenv()
# API_KEY = os.getenv('PEXELS_API_KEY')
API_KEY = 'eK5x0yGyWTU9ChaDciVoBaBnDtrkAJBTXnPYui5UNimFVWNvtI4cWrPU'

UNSPLASH_ACCESS_KEY = '86eEHBr3Dv7VH_izGavj5RXVUsOKsNBH5BFvy-bhUvo'

def parse_response(response):
    slides = response.split('\n\n')
    slides_content = []
    for slide in slides:
        lines = slide.split('\n')
        title_line = lines[0]
        if ': ' in title_line:
            title = title_line.split(': ', 1)[1]  # Extract the title after 'Slide X: '
        else:
            title = title_line
            continue
        content_lines = [line for line in lines[1:] if line != 'Content:']  # Skip line if it is 'Content:'
        content = '\n'.join(content_lines)  # Join the lines to form the content
        # Extract the keyword from the line that starts with 'Keyword:'
        # keyword_line = [line for line in lines if 'Keyword:' or 'Keywords:' in line][0]
        # keyword = keyword_line.split(': ', 1)
        title = title.replace('/n', '')
        content = content.replace('/n', '')
        content = content.replace('- ', '')
        content = content.replace('Content:', '')
        parags = content.split('\n')
        final_content = ''
        for item in parags:
            if item != '':
                final_content += f"{item}\n"
        slides_content.append({'title': title, 'content': final_content, 'keyword': 'key'})
    return slides_content


def search_pexels_images(keyword):
    # num = random.choice([0,1])
    # query = quote_plus(keyword.lower())
    # print("Query:", query) # Debug
    # PEXELS_API_URL = f'https://api.pexels.com/v1/search?query={query}&per_page=1'
    # print("URL:", PEXELS_API_URL) # Debug
    # headers = {
    #     'Authorization': API_KEY
    # }
    # response = requests.get(PEXELS_API_URL, headers=headers)
    # print("Response Status Code:", response.status_code) # Debug
    # print("Response Content:", response.text) # Debug
    # data = json.loads(response.text)
    # if 'photos' in data:
    #     if len(data['photos']) > 0:
    #         return data['photos'][0]['src']['medium']
    # return None
    pu = PyUnsplash(api_key=UNSPLASH_ACCESS_KEY)
    query = quote_plus(keyword.lower())
    photos = pu.photos(type_='random', count=1, featured=True, query=query)
    [photo] = photos.entries
    print(photo.id, photo.link_download)
    response = requests.get(photo.link_download, allow_redirects=True)
    unsplash_image = datetime.datetime.now().timestamp()
    # open(f'GeneratedImgs/{unsplash_image}.jpg', 'wb').write(response.content)
    return response.content


def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(slides_content, template_choice, presentation_title, presenter_name, isImage, numSlides, language):
    template_path = f"Designs/{template_choice}.pptx"

    prs = Presentation(template_path)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title

    n = numSlides + 2
    nums = range(2, n-1)

    selectedNums = [2]
    if n > 3:
        selectedNums = random.choices(nums, k=n//2)

    selectedNums.sort()
    print(selectedNums)

    index = 1

    #add subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"

    if template_choice == 'dark_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color

    elif template_choice == 'bright_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Arial'
                run.font.color.rgb = RGBColor(255, 20, 147)  # RGB for deep pink color

    # add content slides
    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)
        
        if index in selectedNums:
            # fetch image URL from Pixabay based on the slide's title
            image_data = search_pexels_images(presentation_title)
            # print("Image URL:", image_url) #debug
            if image_data is not None:
                # download the image
                # image_data = requests.get(image_url).content
                # load image into BytesIO object
                image_stream = io.BytesIO(image_data)
                # add the image at the specified position
                # slide_width = Inches(20)
                # slide_height = Inches(15)

                image_width = Inches(5.76)  # width of image
                image_height = Inches(5.73)  # height of image

                # left = slide_width - image_width  # calculate left position
                # top = slide_height - image_height - Inches(4)  # calculate top position
                left = Inches(12.6)
                top = Inches(3.63)

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color
            elif placeholder.placeholder_format.type == 7:  # Content
                text = slide_content['content']
                sentences = []
                if language == 'Chinese':
                    sentences = [sentence.strip() for sentence in text.split('。') if sentence.strip()]
                else:
                    sentences = [sentence.strip() for sentence in text.split('.') if sentence.strip()]
                # Join the non-empty sentences back into a string
                result = '\n '.join(sentences)
                placeholder.text = result
                if index in selectedNums:
                    placeholder.left = Inches(1.5)
                    placeholder.top = Inches(3.44)
                    placeholder.width = Inches(10.92)
                    placeholder.height = Inches(6.06)
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        if template_choice == 'dark_modern':
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for white color
        index = index + 1

    # add credits slide
    # slide = prs.slides.add_slide(content_slide_layout)
    # if template_choice == 'dark_modern':
    #     for placeholder in slide.placeholders:
    #         if placeholder.placeholder_format.type == 1:  # Title
    #             placeholder.text = "Credits"
    #             for paragraph in placeholder.text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     run.font.name = 'Times New Roman'
    #                     run.font.color.rgb = RGBColor(255, 165, 0)
    #         elif placeholder.placeholder_format.type == 7:  # Content
    #             placeholder.text = "Images provided by Pexels: https://www.pexels.com"
    #             for paragraph in placeholder.text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     run.font.name = 'Times New Roman'
    #                     run.font.color.rgb = RGBColor(255, 255, 255)

    # elif template_choice == 'bright_modern':
    #     for placeholder in slide.placeholders:
    #         if placeholder.placeholder_format.type == 1:  # Title
    #             placeholder.text = "Credits"
    #             for paragraph in placeholder.text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     run.font.name = 'Arial'
    #                     run.font.color.rgb = RGBColor(255, 20, 147)
    #         elif placeholder.placeholder_format.type == 7:  # Content
    #             placeholder.text = "Images provided by Pexels: https://www.pexels.com"
    #             for paragraph in placeholder.text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     run.font.name = 'Arial'
    #                     run.font.color.rgb = RGBColor(0, 0, 0)

    # else:
    #     for placeholder in slide.placeholders:
    #         if placeholder.placeholder_format.type == 1:  # Title
    #             placeholder.text = "Credits"
    #             for paragraph in placeholder.text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     run.font.name = 'Arial'
    #                     run.font.color.rgb = RGBColor(0, 0, 0)
    #         elif placeholder.placeholder_format.type == 7:  # Content
    #             placeholder.text = "Images provided by Pexels: https://www.pexels.com"
    #             for paragraph in placeholder.text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     run.font.name = 'Arial'
    #                     run.font.color.rgb = RGBColor(0, 0, 0)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)
    timestamp = datetime.datetime.now().timestamp()
    save_path = f"../node/public/GeneratedPresentations/{presenter_name}_{timestamp}.pptx"
    # Save the presentation
    prs.save(save_path)
    return f"GeneratedPresentations/{presenter_name}_{timestamp}"
