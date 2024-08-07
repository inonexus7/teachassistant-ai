#!/usr/bin/env python3
from gptutils import aicomplete, aicomplete_davinci
from pptx import Presentation
from pptx.util import Pt, Inches
import re
import random
import os
import requests
from googlesearch import search
from spire.presentation import Presentation as pres, FileFormat
from spire.presentation.common import *
from urllib.parse import quote_plus
from text_pp import *
from openai import OpenAI
import config

api_key = config.DevelopmentConfig.OPENAI_KEY
client = OpenAI(api_key=api_key)

system_prompt = """You are an assistant that gives the idea for PowerPoint presentations. When answering, give the user the summarized content for each slide based on the number of slide.
                    And the format of the answer must be Slide X(the number of the slide): {title of the content} \n Content: \n content with some bullet points.
                    """

if not os.path.exists("GeneratedPresentations"):
    os.makedirs("GeneratedPresentations")
if not os.path.exists("Cache"):
    os.makedirs("Cache")


def create_ppt_text(title, description, number_of_slides, language):
    if number_of_slides > 20:
        number_of_slides = 20
    user_message = f"I want you to come up with the idea for the PowerPoint. The number of slides is {number_of_slides}, Remeber it is important to generate {number_of_slides} of slides. " \
                       f"The title of content for each slide must be unique," \
                       f"Provide one example about title of each slide, the example should have 3 sentences  with at most 100 words for each paragraph the most important aspects of following topic: {description} for each slide. All content should be written by an language - {language}"
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_message}
    ]
    result = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages
    )
    response = result.choices[0].message.content
    return response


def get_image_urls(query, num_results=5):
    image_urls = []

    # Perform a Google image search
    extra_params = {'imgType': 'photo'}
    for result in search(query, num=num_results, stop=num_results, extra_params=extra_params):
        print(result)
        if result.endswith('.jpg') or result.endswith('.png'):
            image_urls.append(result)

    return image_urls

def create_pptx_func(text_file, presentation_title, presenter_name, insert_image, isImage, numSlides, language):
    # prs = Presentation(f"Designs/Design-3.pptx")
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    totalContent = ""
    with open(text_file, "r", encoding="utf-8") as f:
        imageHeader = ''
        imageContent = ''
        for _, line in enumerate(f):
            totalContent += line
            # if line.startswith("#Title:"):
            #     header = line.replace("#Title:", "").strip()
            #     slide = prs.slides.add_slide(prs.slide_layouts[0])
            #     title = slide.shapes.title
            #     title.text = header
            #     imageHeader = header
            #     body_shape = slide.shapes.placeholders[1]
            #     continue
            # elif line.startswith("#Slide:"):
            #     if slide_count > 0:
            #         slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
            #         title = slide.shapes.title
            #         title.text = header
            #         body_shape = slide.shapes.placeholders[slide_placeholder_index]
            #         tf = body_shape.text_frame
            #         tf.text = content
            #         imageContent = content
            #         if slide_layout_index == 8:
            #             # fetch image URL from Pixabay based on the slide's title
            #             keyword = f"{imageHeader}, {imageContent}"
            #             image_url = search_pexels_images(keyword)
            #             print("Image URL:", image_url) #debug
            #             if image_url is not None:
            #                 # download the image
            #                 image_data = requests.get(image_url).content
            #                 # load image into BytesIO object
            #                 image_stream = io.BytesIO(image_data)
            #                 # add the image at the specified position

            #                 height = Inches(4.5)  # width of image
            #                 width = Inches(5.55)  # height of image

            #                 left = Inches(7.4)  # calculate left position
            #                 top = Inches(1.5)  # calculate top position

            #                 slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
            #     content = ""
            #     slide_count += 1
            #     slide_layout_index = last_slide_layout_index
            #     layout_indices = [1, 8]
            #     while slide_layout_index == last_slide_layout_index:
            #         if firsttime == True:
            #             slide_layout_index = 1
            #             slide_placeholder_index = 1
            #             firsttime = False
            #             break
            #         slide_layout_index = random.choice(
            #             layout_indices
            #         )  # Select random slide index
            #         if slide_layout_index == 8:
            #             slide_placeholder_index = 2
            #         else:
            #             slide_placeholder_index = 1
            #     last_slide_layout_index = slide_layout_index
            #     continue

            # elif line.startswith("#Header:"):
            #     header = line.replace("#Header:", "").strip()
            #     continue

            # elif line.startswith("#Content:"):
            #     content = line.replace("#Content:", "").strip()
            #     next_line = f.readline().strip()
            #     while next_line and not next_line.startswith("#"):
            #         content += "\n" + next_line
            #         next_line = f.readline().strip()
            #     continue
    print("=========================================")
    print(totalContent)
    slides_content = parse_response(totalContent)
    template_choice = insert_image
    file_path = create_ppt(slides_content, template_choice, presentation_title, presenter_name, isImage, numSlides, language)
    # timestamp = datetime.datetime.now().timestamp()
    # save_path = f"../node/public/GeneratedPresentations/{userid}_{timestamp}.pptx"
    # prs.save(save_path)
    print('save path: ', file_path)
    
    # saveAsPdf = pres()
    # saveAsPdf.LoadFromFile(save_path)
    # saveAsPdf.SaveToFile(f"../node/public/GeneratedPresentations/{userid}_{ppt_name}.pdf", FileFormat.PDF)
    # saveAsPdf.Dispose()
    # file_path = f"GeneratedPresentations/{userid}_{timestamp}_{ppt_name}"
    return f"{file_path}"


def get_presentation(
    description, title, presenter, number_of_slides, insert_image, isImage, language="English"
):
    # if "," in description:
    #     l1=[_.strip() for _ in description.split(",")]
    # else:
    #     l1=[description]
    # print(l1)
    # x=pdf2final_list.process(l1)
    # print("\n\n", x)
    # file_path = text2ppt.presentate(x)
    description_string = description
    description_string = re.sub(r"[^\w\s.\-\(\)]", "", description_string)
    description_string = description_string.replace("\n", "")
    with open(f"Cache/{description_string}.txt", "w", encoding="utf-8") as f:
        print("creating pptx...")
        f.write(
            create_ppt_text(
                title, description, number_of_slides, language
            )
        )
    file_path = create_pptx_func(
        f"Cache/{description_string}.txt", title, presenter, insert_image, isImage, number_of_slides, language
    )
    print(file_path)
    return str(file_path)
